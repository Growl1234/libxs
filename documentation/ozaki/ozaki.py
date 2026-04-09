#!/usr/bin/env python3
###############################################################################
# Copyright (c) Intel Corporation - All rights reserved.                      #
#                                                                             #
# For information on the license, see the LICENSE file.                       #
# SPDX-License-Identifier: BSD-3-Clause                                       #
###############################################################################
import os
import shutil
import subprocess
import sys

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn


def run_pandoc(args):
    pandoc = shutil.which("pandoc")
    if not pandoc:
        print("ERROR: missing prerequisites!", file=sys.stderr)
        sys.exit(1)
    subprocess.check_call([pandoc] + args)


def autofit(pptx_path):
    """Add normAutofit to every text frame that lacks an autofit setting."""
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                body_pr = shape.text_frame._txBody.find(qn("a:bodyPr"))
                if body_pr is None:
                    continue
                # skip if any autofit element is already present
                if (
                    body_pr.find(qn("a:normAutofit")) is not None
                    or body_pr.find(qn("a:spAutoFit")) is not None
                    or body_pr.find(qn("a:noAutofit")) is not None
                ):
                    continue
                etree.SubElement(body_pr, qn("a:normAutofit"))
    prs.save(pptx_path)


def resave(pptx_path):
    """Toggle auto-fit on each shape via PowerPoint COM to force reflow."""
    winpath = subprocess.check_output(
        ["wslpath", "-w", os.path.abspath(pptx_path)],
        text=True,
    ).strip()
    # ppAutoSizeNone=0, ppAutoSizeTextToFitShape=2
    subprocess.check_call(
        [
            "powershell.exe",
            "-NoProfile",
            "-Command",
            (
                f'$pp = New-Object -ComObject PowerPoint.Application;'
                f'$pres = $pp.Presentations.Open("{winpath}");'
                f'foreach ($slide in $pres.Slides) {{'
                f'  foreach ($shape in $slide.Shapes) {{'
                f'    if ($shape.HasTextFrame) {{'
                f'      $shape.TextFrame2.AutoSize = 0;'
                f'      $shape.TextFrame2.AutoSize = 2'
                f'    }}'
                f'  }}'
                f'}}'
                f'$pres.Save();'
                f'$pres.Close();'
                f'$pp.Quit();'
                f'[System.Runtime.InteropServices.Marshal]'
                f'::ReleaseComObject($pp) | Out-Null'
            ),
        ]
    )


def main():
    here = os.path.dirname(os.path.abspath(__file__))
    output = os.path.join(here, "ozaki.pptx")
    source = os.path.join(here, "index.md")

    run_pandoc([source, "-o", output])
    autofit(output)
    resave(output)


if __name__ == "__main__":
    main()
